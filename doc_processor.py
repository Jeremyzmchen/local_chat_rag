"""
document processor
    1. read documents
    2. clean documents
    3. split documents
"""

import os
from dataclasses import dataclass, field
import time
import re
from pathlib import Path
from io import StringIO
from typing import List, Optional
from pdfminer.high_level import extract_text_to_fp
from langchain_text_splitters import RecursiveCharacterTextSplitter
import hashlib
import logging

logger = logging.getLogger(__name__)


@dataclass
class Chunk:
    """
    A chunk of text, with metadata(id, source, timestamp, etc.).
    """

    chunk_id: str           # {doc_id}_{chunk_index}
    doc_id: str             # belongs to which document
    content: str  
    source: str             # document source / name of the document      
    file_type: str  
    chunk_index: int        # chunk index in the document
    total_chunks: int
    char_count: int = 0     # number of characters in the chunk
    created_at: float = field(default_factory=time.time)

    def __post_init__(self):
        # dataclass will not calculate the length of the content when creating the object
        self.char_count = len(self.content)

    def to_dict(self):
        """
        Convert the Chunk object to a dictionary for FAISS mapping.
        """

        return {
            "chunk_id": self.chunk_id,
            "doc_id": self.doc_id,
            # "content": self.content, # store in faiss_contents_map
            "source": self.source,
            "file_type": self.file_type,
            "chunk_index": self.chunk_index,
            "total_chunks": self.total_chunks,
            "char_count": self.char_count,
            "created_at": self.created_at
        }
    
    
class TextExtractor:
    """
    Text extractor
    """
    def __init__(self):
        """
        Design pattern: Strategy
        Initialize the text extractor with supported file types and their corresponding extraction methods.
        """

        self._extractors = {
            '.pdf':  self._extract_pdf,
            '.txt':  self._extract_txt,
            '.md':   self._extract_txt, 
            '.docx': self._extract_docx,
            '.xlsx': self._extract_excel,
            '.xls':  self._extract_excel,
            '.pptx': self._extract_pptx,
        }

    def extract(self, filepath: str) -> str:
        """
        Main entry point for text extraction.
        Determine the file type and call the corresponding extraction method.
        """

        suffix = Path(filepath).suffix.lower()
        extractor = self._extractors.get(suffix)

        if extractor is None:
            raise ValueError(f"Unsupported file type: {suffix}")

        logger.info(f"Extracting text from {Path(filepath).name} Type: {suffix}")
        return extractor(filepath)
    
    @staticmethod
    def _extract_pdf(filepath: str) -> str:
        """
        Extract text from PDF file.
        Using StringIO to avoid writing to disk.
        """

        # prepare output space
        output = StringIO()
        with open(filepath, 'rb') as f:
            extract_text_to_fp(f, output)
        return output.getvalue()

    @staticmethod
    def _extract_txt(filepath: str) -> str:
        """
        Extract text from TXT or MD file.
        Provide a fallback for different encoding.
        """

        for encoding in ('utf-8', 'gbk', 'utf-16'):
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Failed to identify the file: {filepath}")
    
    @staticmethod
    def _extract_docx(filepath: str) -> str:
        """
        Extract text from DOCX file.
        Using python-docx to parse the file.
        """

        try:
            from docx import Document
        except ImportError:
            raise ImportError("Processing .docx file requires: pip install python-docx")
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    
    @staticmethod
    def _extract_excel(filepath: str) -> str:
        """
        Extract text from Excel file.
        Using pandas to parse the file.
        """

        try:
            import pandas as pd
        except ImportError:
            raise ImportError("Processing Excel requires: pip install pandas openpyxl")
        xl = pd.ExcelFile(filepath)
        parts = []
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            parts.append(f"## table: {sheet_name}\n{df.to_string(index=False)}")
        return "\n\n".join(parts)
    
    @staticmethod
    def _extract_pptx(filepath: str) -> str:
        """
        Extract text from PPTX file.
        Using python-pptx to parse the file.
        """

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Processing .pptx requires: pip install python-pptx")
        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text for shape in slide.shapes
                     if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"## Page:{i}\n" + "\n".join(texts))
        return "\n\n".join(slides)



class TextCleaner:
    """
    Text cleaner
    Regular expression
    """

    _RE_CONTROL_CHARS = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]')
    _RE_HYPHEN_NEWLINE = re.compile(r'([a-zA-Z])[‐\-]\n([a-zA-Z])')
    _RE_PAGE_NUMBER = re.compile(r'^\s*\d{1,4}\s*$', re.MULTILINE)
    _RE_HEADER_FOOTER = re.compile(
        r'(第\s*\d+\s*页|Page\s+\d+|[-–—]\s*\d+\s*[-–—])',
        re.IGNORECASE
    )
    _RE_EXCESS_NEWLINES = re.compile(r'\n{3,}')
    _RE_INLINE_SPACES = re.compile(r'[^\S\n]+')
    _RE_CN_EN_BOUNDARY_1 = re.compile(r'([\u4e00-\u9fff])([A-Za-z0-9])')   
    _RE_CN_EN_BOUNDARY_2 = re.compile(r'([A-Za-z0-9])([\u4e00-\u9fff])')

    def clean(self, text: str) -> str:
        if not text or not text.strip():
            return ""
        text = self._RE_CONTROL_CHARS.sub('', text)
        # "natu-\nral" → "natural"
        text = self._RE_HYPHEN_NEWLINE.sub(r'\1\2', text)
        text = self._RE_PAGE_NUMBER.sub('', text)
        text = self._RE_HEADER_FOOTER.sub('', text)
        text = self._RE_EXCESS_NEWLINES.sub('\n\n', text)
        text = self._RE_INLINE_SPACES.sub(' ', text)
        text = self._RE_CN_EN_BOUNDARY_1.sub(r'\1 \2', text)
        text = self._RE_CN_EN_BOUNDARY_2.sub(r'\1 \2', text)
        return text.strip()
    

class Chunker:
    """
    Text chunker
        - RecursiveCharacterTextSplitter
            chunk according to priority
        - chunk_overlap: about 15% of chunk_size
    """

    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 128):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

        # recursive splitter
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=[
                "\n\n",      # paragraph
                "\n",        
                "。", "！", "？",    # chn sentence ending
                ".", "!", "?",       # eng sentence ending
                "；", ";",     
                "，", ",",     
                " ",       
                "",     
            ],
            length_function=len,
        )

    def split(self, text: str) -> List[str]:
        if not text or not text.strip(): return []
        chunks = self._splitter.split_text(text)
        return [chunk for chunk in chunks if len(chunk.strip()) >= 10]
    

class DocProcessor:
    """
    Main entry point for document processing
    Including:
        - Text extractor
        - Text cleaner
        - Text chunker
    """

    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 128):
        self.extractor = TextExtractor()
        self.cleaner = TextCleaner()
        self.chunker = Chunker(chunk_size, chunk_overlap)

    @staticmethod
    def _generate_doc_id(filepath: str) -> str:
        with open(filepath, 'rb') as f:
            content_hash = hashlib.md5(f.read()).hexdigest()[:12]
        # Path.stem: trips the suffix
        stem = Path(filepath).stem
        return f"doc_{stem}_{content_hash}"
        # example: "doc_report_d41d8cd98f00"

    def process(self, filepath: str, doc_id: Optional[str] = None,) -> List[Chunk]:
        """
        Process document
        
        Returns:
            List[Chunk]
        """

        filepath = str(filepath)    # Path or str
        filename = Path(filepath).name
        file_type = Path(filepath).suffix.lstrip('.').lower()

        # if doc_id is None, generate a unique id
        if doc_id is None:
            doc_id = self._generate_doc_id(filepath)

        logger.info(f"Start processing the doc: {filename} (doc_id={doc_id})")

        # ── Step 1: extract raw text ──
        try:
            raw_text = self.extractor.extract(filepath)
        except Exception as e:
            logger.error(f"Failed to extract doc [{filename}]: {e}")
            return []

        if not raw_text or not raw_text.strip():
            logger.warning(f"Doc's content is empty: {filename}")
            return []

        # ── Step 2: clean text ──
        cleaned_text = self.cleaner.clean(raw_text)

        if not cleaned_text:
            logger.warning(f"The final text after clean is empty: {filename}")
            return []

        logger.info(f"Clean completed: Length of raw_text: {len(raw_text)}  → After clean: {len(cleaned_text)} ")

        # ── Step 3: Chunk text ──
        text_chunks = self.chunker.split(cleaned_text)

        if not text_chunks:
            # For example: if pdf only has img or length of text is less than 10
            logger.warning(f"Chunks are empty: {filename}")
            return []

        # ── Step 4: Encapsulate chunks ──
        total = len(text_chunks)
        chunks = []
        for i, content in enumerate(text_chunks):
            chunk = Chunk(
                chunk_id = f"{doc_id}_chunk_{i}",
                doc_id = doc_id,
                content = content,
                source = filename,
                file_type = file_type,
                chunk_index = i,
                total_chunks = total,
            )
            chunks.append(chunk)

        logger.info(f"Process completed: {filename} → {total} chunks")
        return chunks
    
    def process_batch(self, 
                      filepaths: List[str], 
                      existing_doc_ids: Optional[set[str]] = None, 
                      progress_callback=None
    ) -> List[Chunk]:
        """
        Process multiple documents in a batch

        Args:
            filepaths: List[str]
            existing_doc_ids: set[str] this
            progress_callback: Callable[[int, int, str], None]
        
        Returns:
            List[Chunk]
        """
        all_chunks = []
        total_files = len(filepaths)
        existing_doc_ids = existing_doc_ids or set()

        for i, filepath in enumerate(filepaths):
            filename = Path(filepath).name

            # skip existing docs
            doc_id = self._generate_doc_id(filepath)
            if doc_id in existing_doc_ids:
                logger.info(f"Doc ID {doc_id} already exists, skip: {filename}")
                continue

            if progress_callback:
                progress_callback(i, total_files, filename)

            chunks = self.process(filepath, doc_id)
            all_chunks.extend(chunks)

            logger.info(
                f"Batch progress: [{i+1}/{total_files}]: "
                f"{filename} → {len(chunks)} chunks，total: {len(all_chunks)}"
            )

        if progress_callback:
            progress_callback(total_files, total_files, "Completed")

        return all_chunks


